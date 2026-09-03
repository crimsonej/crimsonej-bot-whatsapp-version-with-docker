"""
services/doc_writer.py
======================
Multi-format document generator engine.
Generates Word (.docx), PDF (.pdf), PowerPoint (.pptx), and Excel (.xlsx) files.
"""

from __future__ import annotations

import os
import re
import time
import random
import tempfile
from typing import Any

from core.config import log

def create_document_file(
    doc_format: str,
    title: str,
    content_data: str | list | dict,
    filename: str | None = None
) -> tuple[str, str] | None:
    """
    Unified document creation entry point.
    Returns (file_path, output_filename) on success, or None on failure.
    """
    fmt = (doc_format or "docx").lower().strip(".")
    if fmt in ("word", "doc"):
        fmt = "docx"
    elif fmt in ("powerpoint", "ppt", "slides"):
        fmt = "pptx"
    elif fmt in ("excel", "spreadsheet", "xls", "csv"):
        fmt = "xlsx"

    safe_title = re.sub(r"[^\w\s-]", "", title or "document").strip().replace(" ", "_")[:40] or "document"
    out_filename = filename or f"{safe_title}.{fmt}"
    if not out_filename.endswith(f".{fmt}"):
        out_filename += f".{fmt}"

    temp_dir = "/dev/shm" if os.path.isdir("/dev/shm") else tempfile.gettempdir()
    temp_path = os.path.join(temp_dir, f"gen_{int(time.time()*1000)}_{random.randint(1000, 9999)}_{out_filename}")

    try:
        if fmt == "docx":
            path = _build_docx(title, content_data, temp_path)
        elif fmt == "pdf":
            path = _build_pdf(title, content_data, temp_path)
        elif fmt == "pptx":
            path = _build_pptx(title, content_data, temp_path)
        elif fmt == "xlsx":
            path = _build_excel(title, content_data, temp_path)
        else:
            path = _build_docx(title, content_data, temp_path)

        if path and os.path.exists(path) and os.path.getsize(path) > 100:
            log.info("[DocWriter] Successfully generated %s (%d bytes)", out_filename, os.path.getsize(path))
            return path, out_filename
        return None
    except Exception as e:
        log.error("[DocWriter] Failed to generate %s document: %s", fmt, e)
        return None

def _parse_content_sections(content_data: str | list | dict) -> list[dict[str, Any]]:
    """Parse raw text/data into structured sections (headings, paragraphs, bullets, tables)."""
    sections = []
    if isinstance(content_data, list):
        for item in content_data:
            if isinstance(item, dict):
                sections.append(item)
            elif isinstance(item, str):
                sections.append({"type": "paragraph", "text": item})
        return sections

    if isinstance(content_data, dict):
        for key, val in content_data.items():
            sections.append({"type": "heading", "text": str(key)})
            if isinstance(val, list):
                sections.append({"type": "table" if val and isinstance(val[0], (list, dict)) else "bullets", "data": val})
            else:
                sections.append({"type": "paragraph", "text": str(val)})
        return sections

    # Plain text parsing by markdown-style lines
    lines = str(content_data or "").splitlines()
    current_table = []

    for line in lines:
        raw = line.strip()
        if not raw:
            continue

        if raw.startswith("|") and raw.endswith("|"):
            # Table row
            cells = [c.strip() for c in raw.strip("|").split("|")]
            # Ignore markdown header separators like "|---|---|"
            if not all(set(c) <= {"-", ":", " "} for c in cells):
                current_table.append(cells)
            continue
        elif current_table:
            sections.append({"type": "table", "data": current_table})
            current_table = []

        if raw.startswith("#"):
            level = len(raw.split()[0]) if raw.startswith("#") else 1
            text = raw.lstrip("#").strip()
            sections.append({"type": "heading", "level": min(level, 3), "text": text})
        elif raw.startswith(("- ", "* ", "• ")):
            sections.append({"type": "bullet", "text": raw.lstrip("-*• ").strip()})
        else:
            sections.append({"type": "paragraph", "text": raw})

    if current_table:
        sections.append({"type": "table", "data": current_table})

    return sections

def _build_docx(title: str, content_data: Any, out_path: str) -> str:
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = docx.Document()

    # Document Title
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run(title)
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    sections = _parse_content_sections(content_data)

    for sec in sections:
        stype = sec.get("type")
        if stype == "heading":
            level = sec.get("level", 1)
            h = doc.add_heading(sec.get("text", ""), level=level)
            h.runs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        elif stype == "bullet":
            doc.add_paragraph(sec.get("text", ""), style="List Bullet")
        elif stype == "paragraph":
            doc.add_paragraph(sec.get("text", ""))
        elif stype == "table":
            tdata = sec.get("data", [])
            if tdata:
                table = doc.add_table(rows=len(tdata), cols=len(tdata[0]))
                table.style = "Table Grid"
                for r_idx, row in enumerate(tdata):
                    for c_idx, val in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = str(val)
                        if r_idx == 0:
                            # Bold headers
                            for p in cell.paragraphs:
                                for r in p.runs:
                                    r.font.bold = True

    doc.save(out_path)
    return out_path

def _build_pdf(title: str, content_data: Any, out_path: str) -> str:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    doc = SimpleDocTemplate(out_path, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=22,
        leading=26,
        textColor=colors.HexColor("#1F497D"),
        alignment=1, # Center
        spaceAfter=20
    )
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading2'],
        fontSize=14,
        leading=18,
        textColor=colors.HexColor("#1F497D"),
        spaceBefore=12,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['BodyText'],
        fontSize=10,
        leading=14,
        spaceAfter=6
    )

    story = [Paragraph(title, title_style), Spacer(1, 10)]
    sections = _parse_content_sections(content_data)

    for sec in sections:
        stype = sec.get("type")
        if stype == "heading":
            story.append(Paragraph(sec.get("text", ""), h2_style))
        elif stype == "bullet":
            story.append(Paragraph(f"• {sec.get('text', '')}", body_style))
        elif stype == "paragraph":
            story.append(Paragraph(sec.get("text", ""), body_style))
        elif stype == "table":
            tdata = sec.get("data", [])
            if tdata:
                table_data = []
                for row in tdata:
                    table_data.append([Paragraph(str(c), body_style) for c in row])

                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#EAECEE")),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.HexColor("#1F497D")),
                    ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#BDC3C7")),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 6),
                    ('TOPPADDING', (0,0), (-1,-1), 6),
                ]))
                story.append(t)
                story.append(Spacer(1, 10))

    doc.build(story)
    return out_path

def _build_pptx(title: str, content_data: Any, out_path: str) -> str:
    import pptx
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = pptx.Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Generated by Crimsonej AI Engine"

    sections = _parse_content_sections(content_data)

    # Group sections into slides (each heading starts a new slide)
    current_slide = None
    current_tf = None

    for sec in sections:
        stype = sec.get("type")
        if stype == "heading" or current_slide is None:
            bullet_slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(bullet_slide_layout)
            slide_title = sec.get("text", title) if stype == "heading" else title
            current_slide.shapes.title.text = slide_title
            current_tf = current_slide.placeholders[1].text_frame
            current_tf.text = ""
        else:
            if stype == "bullet" or stype == "paragraph":
                p = current_tf.add_paragraph()
                p.text = sec.get("text", "")
                if stype == "bullet":
                    p.level = 0
            elif stype == "table":
                tdata = sec.get("data", [])
                if tdata:
                    rows = len(tdata)
                    cols = len(tdata[0])
                    top = Inches(2.5)
                    left = Inches(1.0)
                    width = Inches(8.0)
                    height = Inches(0.8 * rows)
                    table_shape = current_slide.shapes.add_table(rows, cols, left, top, width, height)
                    table = table_shape.table
                    for r_idx, row in enumerate(tdata):
                        for c_idx, val in enumerate(row):
                            table.cell(r_idx, c_idx).text = str(val)

    prs.save(out_path)
    return out_path

def _build_excel(title: str, content_data: Any, out_path: str) -> str:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:30] if title else "Sheet1"

    # Header title banner
    ws.merge_cells("A1:E1")
    title_cell = ws["A1"]
    title_cell.value = title
    title_cell.font = Font(name="Calibri", size=16, bold=True, color="FFFFFF")
    title_cell.fill = PatternFill(start_color="1F497D", end_color="1F497D", fill_type="solid")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 35

    sections = _parse_content_sections(content_data)
    row_cursor = 3

    for sec in sections:
        stype = sec.get("type")
        if stype == "heading":
            ws.cell(row=row_cursor, column=1, value=sec.get("text", "")).font = Font(name="Calibri", size=13, bold=True, color="1F497D")
            row_cursor += 1
        elif stype == "paragraph" or stype == "bullet":
            ws.cell(row=row_cursor, column=1, value=sec.get("text", "")).font = Font(name="Calibri", size=11)
            row_cursor += 1
        elif stype == "table":
            tdata = sec.get("data", [])
            if tdata:
                header_fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
                thin_border = Border(
                    left=Side(style='thin', color='BFBFBF'),
                    right=Side(style='thin', color='BFBFBF'),
                    top=Side(style='thin', color='BFBFBF'),
                    bottom=Side(style='thin', color='BFBFBF')
                )

                for r_idx, row_vals in enumerate(tdata):
                    for c_idx, val in enumerate(row_vals):
                        cell = ws.cell(row=row_cursor + r_idx, column=c_idx + 1, value=val)
                        cell.border = thin_border
                        if r_idx == 0:
                            cell.font = Font(name="Calibri", size=11, bold=True, color="1F497D")
                            cell.fill = header_fill

                row_cursor += len(tdata) + 1

    # Autofit column widths
    for col in ws.columns:
        max_len = max(len(str(cell.value or '')) for cell in col)
        col_letter = get_column_letter(col[0].column)
        ws.column_dimensions[col_letter].width = max(max_len + 3, 12)

    wb.save(out_path)
    return out_path
