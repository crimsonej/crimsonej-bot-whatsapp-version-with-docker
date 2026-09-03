"""
services/doc_reader.py
======================
Multi-format document parsing and text extraction engine.
Supports PDF, DOCX, PPTX, XLSX/CSV, and Text/Code files.
"""

from __future__ import annotations

import io
import os
import re
import base64
from typing import Any

from core.config import log

def extract_document_text(
    doc_bytes_or_b64: bytes | str,
    filename: str = "document",
    mimetype: str = ""
) -> dict[str, Any]:
    """
    Extract readable text, structural context, and metadata from various document formats.
    Returns:
    {
        "ok": bool,
        "filename": str,
        "format": str,
        "text": str,
        "summary": str,
        "metadata": dict
    }
    """
    # Convert base64 string to bytes if needed
    if isinstance(doc_bytes_or_b64, str):
        try:
            # Handle potential data-url header e.g. "data:application/pdf;base64,..."
            if "," in doc_bytes_or_b64 and "base64" in doc_bytes_or_b64[:50]:
                doc_bytes_or_b64 = doc_bytes_or_b64.split(",", 1)[1]
            raw_bytes = base64.b64decode(doc_bytes_or_b64)
        except Exception as e:
            log.warning("[DocReader] base64 decode failed for %s: %s", filename, e)
            return {"ok": False, "filename": filename, "format": "unknown", "text": "", "summary": "Failed to decode base64 document", "metadata": {}}
    else:
        raw_bytes = doc_bytes_or_b64

    ext = os.path.splitext(filename)[1].lower().strip(".")
    mtype = (mimetype or "").lower()

    # Route based on file extension / MIME type
    if ext == "pdf" or "pdf" in mtype:
        return _parse_pdf(raw_bytes, filename)
    elif ext in ("docx", "doc") or "wordprocessingml" in mtype:
        return _parse_docx(raw_bytes, filename)
    elif ext in ("pptx", "ppt") or "presentationml" in mtype:
        return _parse_pptx(raw_bytes, filename)
    elif ext in ("xlsx", "xls", "csv") or "spreadsheetml" in mtype or "csv" in mtype:
        return _parse_excel_or_csv(raw_bytes, filename, ext)
    else:
        # Fallback to plain text parsing
        return _parse_text(raw_bytes, filename)

def _parse_pdf(data: bytes, filename: str) -> dict[str, Any]:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        num_pages = len(reader.pages)
        pages_text = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                pages_text.append(f"--- Page {i+1} ---\n{txt.strip()}")

        full_text = "\n\n".join(pages_text)
        meta = reader.metadata or {}
        info = {
            "num_pages": num_pages,
            "title": getattr(meta, "title", None) or "",
            "author": getattr(meta, "author", None) or "",
        }
        summary = f"PDF document '{filename}' ({num_pages} pages, {len(full_text)} characters)"
        return {
            "ok": True,
            "filename": filename,
            "format": "pdf",
            "text": full_text,
            "summary": summary,
            "metadata": info
        }
    except Exception as e:
        log.warning("[DocReader] PDF parse error for %s: %s", filename, e)
        return _parse_text(data, filename)

def _parse_docx(data: bytes, filename: str) -> dict[str, Any]:
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        elements = []

        # Read paragraphs & headings
        for p in doc.paragraphs:
            if not p.text.strip():
                continue
            if p.style and p.style.name.startswith("Heading"):
                elements.append(f"\n### {p.text.strip()}")
            else:
                elements.append(p.text.strip())

        # Read tables
        for t_idx, table in enumerate(doc.tables):
            elements.append(f"\n[Table {t_idx+1}]")
            for row in table.rows:
                row_vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                elements.append(" | ".join(row_vals))

        full_text = "\n".join(elements)
        summary = f"Word document '{filename}' ({len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables)"
        return {
            "ok": True,
            "filename": filename,
            "format": "docx",
            "text": full_text,
            "summary": summary,
            "metadata": {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
        }
    except Exception as e:
        log.warning("[DocReader] DOCX parse error for %s: %s", filename, e)
        return _parse_text(data, filename)

def _parse_pptx(data: bytes, filename: str) -> dict[str, Any]:
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(data))
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_parts = [f"--- Slide {i+1} ---"]

            # Extract title if present
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_parts.append(f"Title: {slide.shapes.title.text.strip()}")

            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_parts.append(f"• {paragraph.text.strip()}")

                # Extract table text if present
                if shape.has_table:
                    slide_parts.append("[Slide Table]")
                    for row in shape.table.rows:
                        row_cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                        slide_parts.append(" | ".join(row_cells))

            # Extract notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"[Notes]: {notes}")

            slides_text.append("\n".join(slide_parts))

        full_text = "\n\n".join(slides_text)
        summary = f"PowerPoint presentation '{filename}' ({len(prs.slides)} slides)"
        return {
            "ok": True,
            "filename": filename,
            "format": "pptx",
            "text": full_text,
            "summary": summary,
            "metadata": {"num_slides": len(prs.slides)}
        }
    except Exception as e:
        log.warning("[DocReader] PPTX parse error for %s: %s", filename, e)
        return _parse_text(data, filename)

def _parse_excel_or_csv(data: bytes, filename: str, ext: str) -> dict[str, Any]:
    try:
        import pandas as pd
        sheets_out = []

        if ext == "csv" or filename.lower().endswith(".csv"):
            df = pd.read_csv(io.BytesIO(data))
            sheet_str = f"--- CSV Data: {filename} ---\nColumns: {list(df.columns)}\nRows: {len(df)}\n\n"
            sheet_str += df.head(50).to_markdown(index=False)
            full_text = sheet_str
            meta = {"sheets": ["CSV"], "total_rows": len(df)}
        else:
            excel_file = pd.ExcelFile(io.BytesIO(data))
            for sheet in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet)
                s_str = f"--- Sheet: {sheet} ---\nColumns: {list(df.columns)}\nRows: {len(df)}\n\n"
                s_str += df.head(30).to_markdown(index=False)
                sheets_out.append(s_str)
            full_text = "\n\n".join(sheets_out)
            meta = {"sheets": excel_file.sheet_names}

        summary = f"Spreadsheet '{filename}' ({len(meta.get('sheets', []))} sheets)"
        return {
            "ok": True,
            "filename": filename,
            "format": "excel",
            "text": full_text,
            "summary": summary,
            "metadata": meta
        }
    except Exception as e:
        log.warning("[DocReader] Excel/CSV parse error for %s: %s", filename, e)
        return _parse_text(data, filename)

def _parse_text(data: bytes, filename: str) -> dict[str, Any]:
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception:
        text = data.decode("latin-1", errors="replace")

    # Sanitize unprintable binary null bytes
    clean_text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
    summary = f"Text document '{filename}' ({len(clean_text)} characters)"
    return {
        "ok": True,
        "filename": filename,
        "format": "text",
        "text": clean_text,
        "summary": summary,
        "metadata": {"char_count": len(clean_text)}
    }
